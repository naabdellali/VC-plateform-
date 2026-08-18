import io

from pptx import Presentation
from pptx.util import Inches

from app.services.deck_parser import parse_deck


def _build_sample_pptx() -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Market Opportunity"
    body = slide.placeholders[1]
    body.text_frame.text = "TAM: EUR 10bn\nSAM: EUR 2bn\nSOM: EUR 200m"

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Traction"
    body2 = slide2.placeholders[1]
    body2.text_frame.text = "Current MRR: EUR 85k\nGrowth: 15% MoM"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_pptx_extracts_slide_text_and_titles():
    deck_bytes = _build_sample_pptx()
    parsed = parse_deck("pitch.pptx", deck_bytes)

    assert len(parsed.slides) == 2
    assert parsed.slides[0].title == "Market Opportunity"
    assert "TAM" in parsed.slides[0].text
    assert "SAM" in parsed.slides[0].text
    assert parsed.slides[1].title == "Traction"
    assert "MRR" in parsed.slides[1].text

    raw = parsed.raw_text
    assert "Slide 1: Market Opportunity" in raw
    assert "Slide 2: Traction" in raw


def test_parse_deck_rejects_unsupported_format():
    import pytest

    with pytest.raises(ValueError):
        parse_deck("deck.key", b"not a real file")


def test_slides_json_serializable():
    import json

    deck_bytes = _build_sample_pptx()
    parsed = parse_deck("pitch.pptx", deck_bytes)
    # must round-trip through json.dumps without error - this is what gets
    # stored in Deck.slides_json
    serialized = json.dumps(parsed.to_slides_json())
    assert "Market Opportunity" in serialized
