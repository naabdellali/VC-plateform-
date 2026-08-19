"""
Phase 1 demonstration - run against a dense synthetic deck through the REAL
upload endpoint (not a unit test), in mock mode (no ANTHROPIC_API_KEY in this
sandbox), and print exactly what the OLD pipeline (extract_claims -> flat
Deck.extracted_claims_json) captured vs. what the NEW canonical deal
representation (Number + Claim tables, Pass A-E) captures - including the
"financials"/"other" categories that used to disappear end-to-end because no
reasoning module read them.

Run with:
    python scripts/phase1_demo.py

Honest caveat printed at the end: this sandbox has no live ANTHROPIC_API_KEY,
so every LLM-shaped pass below (classification, structured fields, management
claims, assumption decomposition) runs on its deterministic MOCK fallback -
materially inferior recall to a live Claude call. Pass A (number recognition)
is real, not mocked - it's plain regex, identical in mock and live mode.
"""
import io
import os
import sys
import tempfile

os.environ["DATABASE_URL"] = f"sqlite:///{tempfile.mktemp(suffix='.db')}"
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation  # noqa: E402
from starlette.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402
from app.db import SessionLocal  # noqa: E402
from app.models import Company, Deck, Number, Claim  # noqa: E402


SLIDES = [
    ("Company", "Acme SaaS SAS, siege social a Paris.\nJean Dupont, CEO et co-fondateur, ex-Google.\nMarie Curie, CTO et co-fondatrice, ex-Stripe."),
    ("Problem & Solution", "Le probleme: la gestion des notes de frais est manuelle et lente pour les PME.\nNotre solution: une plateforme SaaS d'automatisation des notes de frais."),
    ("Traction", "MRR actuel: 90K EUR en decembre 2025, en hausse de 120% sur un an.\n35 clients payants a ce jour. ARR de 2M EUR."),
    ("Revenue by quarter", "Q1 2025: 400K EUR de revenue. Q2 2025: 600K EUR de revenue.\nQ3 2025: 850K EUR de revenue. Q4 2025: 1.1M EUR de revenue."),
    ("Market", "TAM estime a 64 milliards de dollars pour ce secteur.\nSAM: 7.4bn. SOM: 740M."),
    ("Competition", "Nous sommes leader du marche en France, aucun concurrent direct de taille comparable.\nTechnologie proprietaire de rapprochement automatique des recus."),
    ("Funding", "Levee de fonds: 3.5M EUR en seed, menee par Example Ventures.\nValorisation post-money: 15M EUR."),
    ("Forecast", "Objectif: ARR de 8M EUR d'ici fin 2027, porte par l'expansion sur le marche americain."),
    ("Pricing", "Pricing: 49 EUR par mois par utilisateur."),
]


def build_dense_deck() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title, body in SLIDES:
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        s.placeholders[1].text_frame.text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def section(title: str) -> None:
    print(f"\n{'=' * 78}\n{title}\n{'=' * 78}")


def main() -> None:
    from app.db import init_db
    init_db()
    client = TestClient(app)

    company_resp = client.post("/companies", json={"name": "Acme SaaS"})
    company_resp.raise_for_status()
    company_id = company_resp.json()["id"]

    deck_bytes = build_dense_deck()
    upload_resp = client.post(
        f"/companies/{company_id}/deck",
        files={"file": ("acme_dense_deck.pptx", deck_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    upload_resp.raise_for_status()

    db = SessionLocal()
    try:
        company = db.query(Company).filter_by(id=company_id).one()
        deck = db.query(Deck).filter_by(company_id=company_id).order_by(Deck.uploaded_at.desc()).first()

        section("OLD PIPELINE - Deck.extracted_claims_json (extract_claims, flat, mock mode)")
        old_claims = deck.extracted_claims_json or []
        if not old_claims:
            print("(empty)")
        old_categories = {}
        for c in old_claims:
            old_categories.setdefault(c.get("category"), []).append(c)
        for cat, items in old_categories.items():
            print(f"\n[{cat}] ({len(items)})")
            for it in items:
                print(f"  - {it.get('claim')!r}  value={it.get('value')!r}")
        print("\n--- What happens to 'financials' / 'other' under the OLD pipeline ---")
        print("These categories are extracted into the JSON blob above (if the mock/live")
        print("extractor produces them at all) but ZERO reasoning modules ever read")
        print("category=='financials' or category=='other' - see app/services/reasoning/*")
        print("and claim_taxonomy.py's CLAIM_TYPE_TO_MODULES for the audit trail. They are")
        print("stored once, then never surfaced again anywhere in the product.")

        section("NEW PIPELINE - canonical deal representation (Number + Claim tables)")

        numbers = db.query(Number).filter_by(company_id=company_id).order_by(Number.slide_reference).all()
        print(f"\n--- Pass A+B: Numbers ({len(numbers)} extracted) ---")
        print("value | unit | currency | period | semantic_category (confidence) | slide | raw_text")
        for n in numbers:
            print(
                f"  {n.value!r:>14} | {n.unit or '-':>6} | {n.currency or '-':>3} | {n.period or '-':>10} | "
                f"{n.semantic_category.value:<14}({n.semantic_confidence or '-'}) | slide {n.slide_reference or '-':>2} | "
                f"{n.raw_text!r}"
            )

        claims = db.query(Claim).filter_by(company_id=company_id).order_by(Claim.slide_reference).all()
        facts = [c for c in claims if c.kind.value == "fact"]
        mgmt = [c for c in claims if c.kind.value == "company_claim"]
        assumptions = [c for c in claims if c.kind.value == "assumption"]

        print(f"\n--- Pass C: structured facts ({len(facts)}) - descriptive, kind=fact ---")
        for c in facts:
            print(f"  [{c.claim_type}] slide {c.slide_reference or '-'}: {c.text!r}")

        print(f"\n--- Pass D: management claims ({len(mgmt)}) - assertions, kind=company_claim, unverified ---")
        for c in mgmt:
            print(f"  [{c.claim_type}] slide {c.slide_reference or '-'}: {c.text!r}")
            print(f"      required_evidence:    {c.required_evidence}")
            print(f"      potential_challenge:  {c.potential_challenge}")

        print(f"\n--- Pass E: decomposed assumptions ({len(assumptions)}) - underpinning the forecast(s) above ---")
        for a in assumptions:
            parent = next((c for c in claims if c.id == a.parent_claim_id), None)
            print(f"  underpins: {(parent.text if parent else '?')!r}")
            print(f"    assumption: {a.text!r}")

        section("RECALL CHECK - information that used to disappear")
        quarterly_numbers = [n for n in numbers if n.slide_reference == "4"]
        print(f"Quarterly revenue numbers captured from Slide 4 ('financials'-shaped data): {len(quarterly_numbers)}")
        for n in quarterly_numbers:
            print(f"  - {n.raw_text!r} (value={n.value}, context={n.context!r})")
        print(
            "\nUnder the OLD pipeline these either never got a dedicated category or landed in "
            "'financials'/'other' and were never read again. Under the NEW pipeline they are "
            "persisted as inspectable Number rows with full provenance (slide, raw text, context), "
            "regardless of whether a reasoning module consumes them yet."
        )

        section("PROVENANCE CHECK - one example traced end to end")
        example = next((n for n in numbers if "35" in n.raw_text), None)
        if example:
            print(f"raw_text:   {example.raw_text!r}")
            print(f"slide:      {example.slide_reference}")
            print(f"context:    {example.context!r}")
            print(f"value/unit: {example.value} {example.unit}")
            print(f"semantic:   {example.semantic_category.value} (confidence={example.semantic_confidence})")
            print(f"deck_id:    {example.deck_id}")
            print(f"company_id: {example.company_id}")

    finally:
        db.close()

    section("HONEST CAVEAT")
    print(
        "This sandbox has no ANTHROPIC_API_KEY configured, so Pass B/C/D/E above ran on their\n"
        "deterministic MOCK fallbacks (regex/keyword heuristics), not a real Claude call. Pass A\n"
        "(number recognition) is NOT mocked - it's plain regex, identical in mock and live mode,\n"
        "and is what's actually doing the heavy lifting in this demo. Live-mode recall (real\n"
        "semantic classification of ambiguous numbers, real structured-field/management-claim\n"
        "extraction, real assumption decomposition) will be materially higher than what's shown\n"
        "here - this demo proves the pipeline runs and the schema captures what was previously\n"
        "lost, not the ceiling of what the LLM passes can do."
    )


if __name__ == "__main__":
    main()
