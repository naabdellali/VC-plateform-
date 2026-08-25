import time
from playwright.sync_api import sync_playwright

BASE = "http://localhost:3000"


def build_sample_pptx(path):
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Market"
    s1.placeholders[1].text_frame.text = "TAM: EUR 8bn\nSAM: EUR 1.5bn"

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Traction"
    s2.placeholders[1].text_frame.text = "Current MRR: EUR 95k"

    s3 = prs.slides.add_slide(layout)
    s3.shapes.title.text = "Team"
    s3.placeholders[1].text_frame.text = "CEO: Jane Doe, ex-Google, 10 years experience"

    s4 = prs.slides.add_slide(layout)
    s4.shapes.title.text = "Competition"
    s4.placeholders[1].text_frame.text = "Main competitor: LegacyCorp"

    prs.save(path)


def main():
    pptx_path = "/tmp/sample_pitch.pptx"
    build_sample_pptx(pptx_path)

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 900})

        print("1. Loading home page...")
        page.goto(BASE, wait_until="networkidle")
        page.screenshot(path="/tmp/shot_1_home.png")

        print("2. Filling new-company form...")
        page.fill('input[placeholder="Acme SaaS"]', "Acme SaaS")
        page.fill('input[placeholder="B2B expense management software"]', "B2B expense management software")
        page.set_input_files('input[type="file"]', pptx_path)
        page.click('button:has-text("Upload & analyze")')

        print("3. Waiting for redirect to company tray...")
        page.wait_for_url("**/company/**", timeout=30000)
        page.wait_for_selector(".tray", timeout=15000)
        time.sleep(1)
        page.screenshot(path="/tmp/shot_2_tray.png", full_page=True)
        print("   URL:", page.url)

        print("4. Opening Market module drill-down...")
        page.click('a:has-text("Market")')
        page.wait_for_selector("text=Conclusion", timeout=15000)
        time.sleep(0.5)
        page.screenshot(path="/tmp/shot_3_market_module.png", full_page=True)

        print("5. Submitting bottom-up recalculation...")
        page.select_option("select", "bottom_up")
        page.fill('textarea', "20% penetration assumed based on category maturity - REVIEW")
        page.click('button:has-text("Recalculate")')
        page.wait_for_selector("text=Platform estimate:", timeout=15000)
        time.sleep(0.5)
        page.screenshot(path="/tmp/shot_4_market_recalculated.png", full_page=True)

        company_url = page.url.split("/module/")[0]
        print("6. Back to tray, opening memo page...")
        page.goto(f"{company_url}/memo", wait_until="networkidle")
        page.click('button:has-text("Generate memo")')
        page.wait_for_selector("text=RECOMMENDATION", timeout=20000)
        time.sleep(0.5)
        page.screenshot(path="/tmp/shot_5_memo.png", full_page=True)

        browser.close()
        print("\nE2E CHECK PASSED - screenshots in /tmp/shot_*.png")


if __name__ == "__main__":
    main()
