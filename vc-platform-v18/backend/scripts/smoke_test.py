"""
End-to-end smoke test against the real FastAPI app (in-process, via
Starlette's TestClient) - exercises the actual HTTP surface a real
frontend would call, in mock mode (no API keys required). Run with:

    python scripts/smoke_test.py
"""
import io
import os
import sys
import tempfile

os.environ["DATABASE_URL"] = f"sqlite:///{tempfile.mktemp(suffix='.db')}"
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation  # noqa: E402
from starlette.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402


def build_sample_deck() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Market"
    s1.placeholders[1].text_frame.text = "TAM: EUR 8bn\nSAM: EUR 1.5bn"

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Traction"
    s2.placeholders[1].text_frame.text = "Current MRR: EUR 95k\nGrowth: strong month over month"

    s3 = prs.slides.add_slide(layout)
    s3.shapes.title.text = "Team"
    s3.placeholders[1].text_frame.text = "CEO: Jane Doe, ex-Google, 10 years experience\nCTO: John Smith, ex-Amazon"

    s4 = prs.slides.add_slide(layout)
    s4.shapes.title.text = "Competition"
    s4.placeholders[1].text_frame.text = "Main competitor: LegacyCorp\nWe are 10x cheaper"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def main():
    print("1. Health check...")
    r = client.get("/health")
    assert r.status_code == 200, r.text
    print("   ", r.json())

    print("2. Create company...")
    r = client.post("/companies", json={"name": "Acme SaaS"})
    assert r.status_code == 200, r.text
    company = r.json()
    company_id = company["id"]
    print("   company_id =", company_id)

    print("3. Upload deck (triggers full auto-analysis pipeline)...")
    deck_bytes = build_sample_deck()
    r = client.post(
        f"/companies/{company_id}/deck",
        files={"file": ("acme_pitch.pptx", deck_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        data={"stage": "seed", "business_model": "saas", "sector": "B2B expense management", "hq_country": "France"},
    )
    assert r.status_code == 200, r.text
    analyze = r.json()
    print("   modules_triggered =", analyze["modules_triggered"])

    print("4. Get tray...")
    r = client.get(f"/companies/{company_id}/tray")
    assert r.status_code == 200, r.text
    for tile in r.json():
        print(f"   [{tile['status']:20s}] {tile['label']}: {tile['headline']}")

    print("5. Get evidence trail for market module...")
    r = client.get(f"/companies/{company_id}/evidence", params={"module": "market"})
    assert r.status_code == 200, r.text
    evidence = r.json()
    print(f"   {len(evidence)} evidence rows for market module")
    for e in evidence[:3]:
        print(f"     - [{e['origin']}/{e['confidence']}] {e['claim'][:80]}")

    print("6. Recalculate market size (human-in-the-loop bottom-up)...")
    r = client.post(
        f"/companies/{company_id}/modules/market/recalculate",
        json={
            "methodology": "bottom_up",
            "inputs": {"num_potential_customers": 60000, "avg_annual_spend_eur": 2400, "realistic_penetration": 0.2},
            "assumptions": ["20% penetration assumed based on category maturity - REVIEW"],
        },
    )
    assert r.status_code == 200, r.text
    print("   ", r.json())

    print("7. Submit MRR series (volatility check)...")
    r = client.post(
        f"/companies/{company_id}/modules/traction/mrr-series",
        json={"monthly_values_eur": [70000, 95000, 68000, 100000, 65000, 110000]},
    )
    assert r.status_code == 200, r.text
    print("   ", r.json())

    print("8. Submit CAC/LTV consistency check...")
    r = client.post(
        f"/companies/{company_id}/modules/traction/cac-ltv-check",
        json={"cac": 4000, "reported_ltv": 600000, "gross_margin": 0.75, "arpa_monthly": 200},
    )
    assert r.status_code == 200, r.text
    print("   ", r.json())

    print("9. Get red flags...")
    r = client.get(f"/companies/{company_id}/red-flags")
    assert r.status_code == 200, r.text
    flags = r.json()
    print(f"   {len(flags)} red flag(s)")
    for f in flags:
        print(f"     - [{f['severity']}] ({f['category']}) {f['explanation'][:100]}")

    print("10. Generate investment memo...")
    r = client.post(f"/companies/{company_id}/memo/generate")
    assert r.status_code == 200, r.text
    memo = r.json()
    print("    recommendation =", memo["recommendation"])
    for s in memo["sections_json"]:
        print(f"    [{s['title']}] {s['body'][:120]}")

    print("11. Refetch memo via GET...")
    r = client.get(f"/companies/{company_id}/memo")
    assert r.status_code == 200, r.text

    print("\nALL SMOKE TESTS PASSED")


if __name__ == "__main__":
    with TestClient(app) as client:
        globals()["client"] = client
        main()
