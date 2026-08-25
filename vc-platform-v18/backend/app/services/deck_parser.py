"""
Structured extraction from an uploaded pitch deck (.pptx or .pdf).

This is a purely mechanical step - no LLM involved - so it is 100%
deterministic and testable: given the same file, you always get the same
slide-by-slide text back. The LLM extraction pass (turning this raw text
into typed claims like "market_size_claim") lives in llm_client.py and
consumes the output of this module.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field


@dataclass
class SlideContent:
    index: int
    title: str
    text: str
    notes: str = ""

    def to_dict(self) -> dict:
        return {"slide": self.index, "title": self.title, "text": self.text, "notes": self.notes}


@dataclass
class ParsedDeck:
    slides: list[SlideContent] = field(default_factory=list)

    @property
    def raw_text(self) -> str:
        parts = []
        for s in self.slides:
            parts.append(f"--- Slide {s.index}: {s.title} ---\n{s.text}")
            if s.notes:
                parts.append(f"[speaker notes] {s.notes}")
        return "\n\n".join(parts)

    def to_slides_json(self) -> list[dict]:
        return [s.to_dict() for s in self.slides]


def parse_pptx(file_bytes: bytes) -> ParsedDeck:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[SlideContent] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                # tables carry a lot of the real numbers in pitch decks
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(c.text for c in row.cells))
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)
                    if not title and shape == slide.shapes.title:
                        title = line
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(SlideContent(index=i, title=title or f"Slide {i}", text="\n".join(texts), notes=notes))
    return ParsedDeck(slides=slides)


def parse_pdf(file_bytes: bytes) -> ParsedDeck:
    import pdfplumber

    slides: list[SlideContent] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            first_line = text.strip().split("\n")[0] if text.strip() else f"Page {i}"
            table_text = []
            for table in page.extract_tables() or []:
                for row in table:
                    table_text.append(" | ".join(c or "" for c in row))
            full_text = text + ("\n" + "\n".join(table_text) if table_text else "")
            slides.append(SlideContent(index=i, title=first_line[:80], text=full_text))
    return ParsedDeck(slides=slides)


def parse_deck(filename: str, file_bytes: bytes) -> ParsedDeck:
    lower = filename.lower()
    if lower.endswith(".pptx"):
        return parse_pptx(file_bytes)
    if lower.endswith(".pdf"):
        return parse_pdf(file_bytes)
    raise ValueError(f"Unsupported deck format: {filename}. Only .pptx and .pdf are supported in V1.")
